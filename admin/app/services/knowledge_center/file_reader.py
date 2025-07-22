import platform
import subprocess
import stat
import os
import uuid

from app.app import app, db
import pymupdf4llm
from openpyxl import load_workbook
import xlrd
from openpyxl.drawing.image import Image as OpenpyxlImage
import shutil
from PIL import Image
import io
import html2text


def pandoc_reader(resource, task_params):
    """
    使用Pandoc读取文档并转换为指定格式。
        输入平台资源对象
        输入任务参数，包含目标格式等信息
        将文档转换为指定格式并返回转换后的资源信息
    :param resource: 平台资源对象
    :param task_params: 任务参数，包含目标格式等信息
    :return:
    """
    # 这里可以使用subprocess调用pandoc命令行工具进行转换
    tgt_format = task_params.get("to_format", "markdown")
    base_path = pick_pandoc_lib()
    target_new_path = resource.resource_path + f".{tgt_format}"
    media_dir = str(os.path.dirname(target_new_path).replace('\\', '/'))
    target_new_media_path = os.path.join(media_dir, 'media')
    if not os.path.exists(target_new_media_path):
        os.makedirs(target_new_media_path, exist_ok=True)
    to_format = {
        "text": "plain",
        "txt": "plain",
    }.get(tgt_format, tgt_format)
    cmd_args = [
        base_path,
        resource.resource_path,
        '-o', target_new_path,
        '--from', resource.resource_format,
        '--to', to_format,
        f'--extract-media={media_dir}',
    ]
    for key, value in task_params.items():
        if key == 'to_format':
            continue
        if value is True:
            cmd_args.append(f'--{key}')
        elif value:
            cmd_args.extend([f'--{key}', str(value)])
    try:
        subprocess.run(cmd_args, check=True, capture_output=True,  text=True)
    except subprocess.CalledProcessError as e:
        return f"转换失败: {str(e)}"
    except Exception as e:
        return f"发生错误: {str(e)}"
    # 检查转换后的文件是否存在
    if not os.path.exists(target_new_path):
        return f"转换后的文件不存在: {target_new_path}"

    update_media_resource_content(target_new_path, target_new_media_path, resource)
    return {
        "id": resource.id,
        "format": tgt_format,
        "path": target_new_path,
    }


def pick_pandoc_lib():
    """
    根据系统环境选择pandoc
    :return:
    """
    system = platform.system()
    machine = platform.machine()
    # 确定基本目录结构
    # 构建版本映射表
    version_map = {
        ('darwin', 'arm64'): 'pandoc-3.7.0.2-arm64-macOS',
        ('darwin', 'x86_64'): 'pandoc-3.7.0.2-x86_64-macOS',
        ('linux', 'x86_64'): 'pandoc-3.7.0.2-linux-amd64',
        ('linux', 'amd64'): 'pandoc-3.7.0.2-linux-amd64',
        ('linux', 'arm64'): 'pandoc-3.7.0.2-linux-arm64',
        ('windows', 'amd64'): 'pandoc-3.7.0.2-windows-x86_64',
        ('windows', 'x86_64'): 'pandoc-3.7.0.2-windows-x86_64'
    }
    target_version = version_map.get((system.lower(), machine.lower()))
    if not target_version:
        raise Exception("Unsupported platform or architecture for Pandoc.")
    pandoc_name = "pandoc" if system.lower() != 'windows' else "pandoc.exe"
    base_path = os.path.join(
        app.config['base_dir'], 'libs', 'pandoc', target_version, 'bin', pandoc_name
    )
    if not os.path.exists(base_path):
        raise Exception(f"Pandoc binary not found at {base_path}.")
    # 如果不是Windows系统，检查并设置执行权限
    if system != 'windows':
        # 获取当前文件权限
        current_mode = os.stat(base_path).st_mode
        # 检查是否有执行权限
        if not (current_mode & stat.S_IXUSR):
            try:
                # 添加用户执行权限
                os.chmod(base_path, current_mode | stat.S_IXUSR | stat.S_IXGRP | stat.S_IXOTH)
            except Exception as e:
                raise Exception(f"Failed to set execute permission for {base_path}: {str(e)}")
    return base_path


def pymupdf_reader(resource, task_params):
    """
    使用PyMuPDF读取PDF文档并转换为Markdown格式。
        输入平台资源对象
        输入任务参数，包含目标格式等信息
        将PDF文档转换为Markdown格式并返回转换后的资源信息
    """
    target_new_path = resource.resource_path + ".markdown"
    target_new_media_path = target_new_path + ".media"
    if not os.path.exists(target_new_media_path):
        os.makedirs(target_new_media_path, exist_ok=True)
    try:
        md_text = pymupdf4llm.to_markdown(
            resource.resource_path, write_images=True,
            image_path=target_new_media_path,
            **task_params
        )
        with open(target_new_path, 'wb') as f:
            f.write(md_text.encode())
    except Exception as e:
        return f"转换失败: {str(e)}"
    # 检查转换后的文件是否存在
    if not os.path.exists(target_new_path):
        return f"转换后的文件不存在: {target_new_path}"
    update_media_resource_content(target_new_path, target_new_media_path, resource)
    return {
        "id": resource.id,
        "format": "markdown",
        "path": target_new_path,
    }


def html2text_reader(resource, task_params):
    """
    使用html2text读取HTML文档并转换为Markdown格式。
    """

    h = html2text.HTML2Text()
    h.ignore_links = False  # 保留链接
    h.ignore_tables = False
    h.escape_all = False
    h.single_line_break = True
    h.strong_mark = '**'  # 强制加粗用 ** 而非 __
    h.emphasis_mark = '*'
    with open(resource.resource_path, 'r', encoding='utf-8') as file:
        html_content = file.read()
    markdown = h.handle(html_content)
    target_new_path = resource.resource_path + ".markdown"
    target_new_media_path = target_new_path + ".media"
    if not os.path.exists(target_new_media_path):
        os.makedirs(target_new_media_path, exist_ok=True)
    try:
        with open(target_new_path, 'w', encoding='utf-8') as f:
            f.write(markdown)
    except Exception as e:
        return f"转换失败: {str(e)}"
    # 检查转换后的文件是否存在
    if not os.path.exists(target_new_path):
        return f"转换后的文件不存在: {target_new_path}"
    return {
        "id": resource.id,
        "format": "markdown",
        "path": target_new_path,
    }


def openpyxl_reader(resource, task_params):

    target_new_path = resource.resource_path + ".markdown"
    target_new_media_path = target_new_path + ".media"
    if not os.path.exists(target_new_media_path):
        os.makedirs(target_new_media_path, exist_ok=True)
    result = {}
    media_references = {}  # 存储媒体文件引用关系
    if resource.resource_format == 'xlsx':
        wb = load_workbook(resource.resource_path, rich_text=True)
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            media_references[sheet_name] = []
            # 处理工作表中的图片
            # 处理工作表中的图片
            for idx, image in enumerate(sheet._images):
                if isinstance(image, OpenpyxlImage):
                    # 生成唯一图片文件名
                    image_filename = f"image_{sheet_name}_{idx}.png"
                    image_path = os.path.join(target_new_media_path, image_filename)

                    try:
                        # 优先尝试使用嵌入的图片数据
                        if hasattr(image, '_data'):
                            # 直接保存嵌入的图片数据
                            with open(image_path, 'wb') as img_file:
                                img_file.write(image._data())
                        elif hasattr(image, 'path'):
                            # 如果是链接到外部文件，尝试读取数据
                            try:
                                with open(image.path, 'rb') as src_file:
                                    with open(image_path, 'wb') as dst_file:
                                        dst_file.write(src_file.read())
                            except FileNotFoundError:
                                # 如果外部文件不存在，尝试从内存中获取
                                if hasattr(image, '_data'):
                                    with open(image_path, 'wb') as img_file:
                                        img_file.write(image._data())
                                else:
                                    continue

                        media_references[sheet_name].append({
                            'type': 'image',
                            'path': image_path,
                            'position': (image.anchor._from.row, image.anchor._from.col) if hasattr(image.anchor,
                                                                                                    '_from') else (0, 0)
                        })
                    except Exception as e:
                        print(f"无法保存图片 {image_filename}: {str(e)}")
                        continue
            content = []
            # 获取表头
            headers = []
            for cell in sheet[1]:
                headers.append(str(cell.value) if cell.value is not None else ' ')

            content.append("\t".join(headers))

            # 获取数据行
            for row in sheet.iter_rows(min_row=2):
                row_values = []
                for cell in row:
                    # 检查单元格是否有超链接
                    hyperlink = cell.hyperlink.target if cell.hyperlink else None
                    value = str(cell.value) if cell.value is not None else ' '
                    # 如果是超链接，添加Markdown格式
                    if hyperlink:
                        value = f"[{value}]({hyperlink})"
                    row_values.append(value)
                content.append("\t".join(row_values))
            # 添加媒体引用标记
            if media_references[sheet_name]:
                content.append("\n\n## 媒体内容")
                for media in media_references[sheet_name]:
                    if media['type'] == 'image':
                        content.append(f"![]({media['path']})")
            result[sheet_name] = "\n".join(content)
    elif resource.resource_format == 'xls':
        wb = xlrd.open_workbook(resource.resource_path)
        for sheet in wb.sheets():
            content = []
            # 处理表头
            headers = [str(sheet.cell_value(0, col)) for col in range(sheet.ncols)]
            content.append("\t".join(headers))

            # 处理数据行
            for row in range(1, sheet.nrows):
                row_values = []
                for col in range(sheet.ncols):
                    value = sheet.cell_value(row, col)
                    row_values.append(str(value) if value != '' else ' ')
                content.append("\t".join(row_values))
            result[sheet.name] = "\n".join(content)
    content = ''
    for sheet_name, sheet_content in result.items():
        content += f"Sheet: {sheet_name}\n{sheet_content}\n\n"
    try:
        with open(target_new_path, 'w', encoding='utf-8') as f:
            f.write(content)
    except Exception as e:
        return f"转换失败: {str(e)}"
    # 检查转换后的文件是否存在
    if not os.path.exists(target_new_path):
        return f"转换后的文件不存在: {target_new_path}"
    update_media_resource_content(target_new_path, target_new_media_path, resource)
    return {
        "id": resource.id,
        "format": "markdown",
        "path": target_new_path,
    }


def pptx_reader(resource, task_params):
    """
    使用python-pptx读取PPTX文档并转换为Markdown格式。
        输入平台资源对象
        输入任务参数，包含目标格式等信息
        将PPTX文档转换为Markdown格式并返回转换后的资源信息
    """
    from pptx import Presentation
    target_new_path = resource.resource_path + ".markdown"
    target_new_media_path = target_new_path + ".media"
    if not os.path.exists(target_new_media_path):
        os.makedirs(target_new_media_path, exist_ok=True)
    prs = Presentation(resource.resource_path)
    markdown_content = []
    for i, slide in enumerate(prs.slides):
        markdown_content.append(f"## Slide {i + 1}\n")
        for shape in slide.shapes:
            # 处理文本
            if hasattr(shape, "text") and shape.text.strip():
                if shape.text.startswith(("- ", "* ")):  # 列表项
                    markdown_content.append(f"{shape.text}\n")
                else:
                    markdown_content.append(f"{shape.text}\n\n")

            # 处理图片（需保存到本地并引用）
            if shape.shape_type == 13:  # 13 是图片类型
                image_path = os.path.join(target_new_media_path, f"slide_{i}_img_{shape.name}.png")
                with open(image_path, "wb") as f:
                    f.write(shape.image.blob)
                markdown_content.append(f"![]({image_path})\n\n")
        markdown_content.append("---\n")
    markdown_text = "\n".join(markdown_content)
    try:
        with open(target_new_path, 'w', encoding='utf-8') as f:
            f.write(markdown_text)
    except Exception as e:
        return f"转换失败: {str(e)}"
    # 检查转换后的文件是否存在
    if not os.path.exists(target_new_path):
        return f"转换后的文件不存在: {target_new_path}"
    update_media_resource_content(target_new_path, target_new_media_path, resource)
    return {
        "id": resource.id,
        "format": "markdown",
        "path": target_new_path,
    }


def text_reader(resource, task_params):
    """
    读取文本文件并返回内容。
        输入平台资源对象
        输入任务参数，包含目标格式等信息
        返回文本内容
    """
    target_new_path = resource.resource_path + ".markdown"
    target_new_media_path = target_new_path + ".media"
    if not os.path.exists(target_new_media_path):
        os.makedirs(target_new_media_path, exist_ok=True)
    try:
        with open(resource.resource_path, 'r', encoding='utf-8') as file:
            content = file.read()
        with open(target_new_path, 'w', encoding='utf-8') as f:
            f.write(content)
    except Exception as e:
        return f"读取失败: {str(e)}"
    # 检查转换后的文件是否存在
    if not os.path.exists(target_new_path):
        return f"转换后的文件不存在: {target_new_path}"
    return {
        "id": resource.id,
        "format": "markdown",
        "path": target_new_path,
    }


def add_media_resource(media_dir, resource, content):
    """
    添加媒体资源到平台资源对象中。
        :param media_dir: 媒体资源路径
        :param resource: 平台资源对象
        :param content: 解析后的内容
        针对目录下的每一个文件，创建一个新的媒体资源对象，并保存到数据库中。
        生成下载链接
        针对解析好的内容，进行链接替换
    """
    from app.models.resource_center.resource_model import ResourceObjectMeta
    from app.utils.oss.oss_client import generate_download_url
    media_files = os.listdir(media_dir)
    if not media_files:
        return content
    for media_file in media_files:
        media_path = os.path.join(media_dir, media_file)
        if os.path.isfile(media_path):
            new_resource = ResourceObjectMeta(
                resource_parent_id=resource.id,
                user_id=resource.user_id,
                resource_name=media_file,
                resource_type="media",
                resource_format=media_file.split('.')[-1],
                resource_path=media_path,
                resource_size_in_MB=os.path.getsize(media_path)/1024/1024,
                resource_source="knowledge_center",
                resource_download_url=generate_download_url(
                    'knowledge_center', media_path, suffix=media_file.split('.')[-1]
                ).json.get("result"),
            )
            db.session.add(new_resource)
    db.session.commit()
    all_commit_media = ResourceObjectMeta.query.filter_by(
        resource_parent_id=resource.id,
        resource_type="media"
    ).all()
    for media in all_commit_media:
        # 替换内容中的链接
        md_path = media.resource_path.replace("\\", "/")
        content = content.replace(
            f'![]({md_path})',
            f'![]({media.resource_download_url})'
        )
    return content


def update_media_resource_content(target_new_path, target_new_media_path, resource):
    """
    更新所有媒体资源的内容，添加下载链接。
    """
    with open(target_new_path, 'r', encoding='utf-8') as f:
        content = f.read()
    new_content = add_media_resource(
        target_new_media_path,
        resource,
        content
    )
    if new_content:
        with open(target_new_path, 'w', encoding='utf-8') as f:
            f.write(new_content)
    return new_content
